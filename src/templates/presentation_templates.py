"""
Professional Presentation Templates for GrandModel
=================================================

Institutional-grade presentation templates for executive summaries, 
performance analysis, risk assessment, and strategy comparison reports.

Features:
- Executive presentation templates
- Investment committee templates
- Risk committee templates
- Board presentation templates
- Client reporting templates
- Regulatory reporting templates
- Performance attribution templates
- Benchmark comparison templates

Author: Agent 6 - Visualization and Reporting System
"""

from typing import Dict, List, Optional, Any, Tuple
from datetime import datetime, timedelta
import pandas as pd
import numpy as np
from pathlib import Path
import json
import logging
from dataclasses import dataclass
from jinja2 import Template
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import plotly.graph_objects as go
import plotly.io as pio
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.backends.backend_pdf import PdfPages
import warnings

warnings.filterwarnings('ignore')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class PresentationConfig:
    """Configuration for presentation templates"""
    company_name: str = "GrandModel Trading System"
    company_logo: Optional[str] = None
    brand_colors: Dict[str, str] = None
    font_family: str = "Calibri"
    template_style: str = "executive"
    
    def __post_init__(self):
        if self.brand_colors is None:
            self.brand_colors = {
                'primary': '#2E86AB',
                'secondary': '#A23B72',
                'accent': '#F18F01',
                'success': '#28A745',
                'warning': '#FFC107',
                'danger': '#DC3545',
                'dark': '#343A40',
                'light': '#F8F9FA'
            }


class PresentationTemplates:
    """
    Professional presentation template system
    """
    
    def __init__(self, config: Optional[PresentationConfig] = None):
        """
        Initialize presentation templates
        
        Args:
            config: Presentation configuration
        """
        self.config = config or PresentationConfig()
        self.templates_dir = Path("/home/QuantNova/GrandModel/src/templates")
        self.output_dir = Path("/home/QuantNova/GrandModel/results/presentations")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Load templates
        self.templates = {}
        self._load_templates()
        
        logger.info("Presentation Templates initialized")
    
    def create_executive_summary_presentation(self,
                                           report_data: Dict[str, Any],
                                           filename: str = "executive_summary") -> str:
        """
        Create executive summary presentation
        
        Args:
            report_data: Report data
            filename: Output filename
            
        Returns:
            Path to generated presentation
        """
        try:
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Set slide size to 16:9
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, report_data)
            
            # Executive summary slide
            self._create_executive_summary_slide(prs, report_data)
            
            # Key metrics slide
            self._create_key_metrics_slide(prs, report_data)
            
            # Performance overview slide
            self._create_performance_overview_slide(prs, report_data)
            
            # Risk assessment slide
            self._create_risk_assessment_slide(prs, report_data)
            
            # Key insights slide
            self._create_key_insights_slide(prs, report_data)
            
            # Recommendations slide
            self._create_recommendations_slide(prs, report_data)
            
            # Appendix slide
            self._create_appendix_slide(prs, report_data)
            
            # Save presentation
            ppt_path = self.output_dir / f"{filename}.pptx"
            prs.save(str(ppt_path))
            
            logger.info(f"Executive summary presentation created: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error creating executive summary presentation: {e}")
            return ""
    
    def create_performance_analysis_presentation(self,
                                              report_data: Dict[str, Any],
                                              filename: str = "performance_analysis") -> str:
        """
        Create performance analysis presentation
        
        Args:
            report_data: Report data
            filename: Output filename
            
        Returns:
            Path to generated presentation
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, report_data, "Performance Analysis")
            
            # Performance summary
            self._create_performance_summary_slide(prs, report_data)
            
            # Returns analysis
            self._create_returns_analysis_slide(prs, report_data)
            
            # Risk-adjusted metrics
            self._create_risk_adjusted_metrics_slide(prs, report_data)
            
            # Drawdown analysis
            self._create_drawdown_analysis_slide(prs, report_data)
            
            # Trade analysis
            self._create_trade_analysis_slide(prs, report_data)
            
            # Benchmark comparison
            self._create_benchmark_comparison_slide(prs, report_data)
            
            # Time-based analysis
            self._create_time_based_analysis_slide(prs, report_data)
            
            # Statistical analysis
            self._create_statistical_analysis_slide(prs, report_data)
            
            # Conclusions and recommendations
            self._create_conclusions_slide(prs, report_data)
            
            # Save presentation
            ppt_path = self.output_dir / f"{filename}.pptx"
            prs.save(str(ppt_path))
            
            logger.info(f"Performance analysis presentation created: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error creating performance analysis presentation: {e}")
            return ""
    
    def create_risk_assessment_presentation(self,
                                         report_data: Dict[str, Any],
                                         filename: str = "risk_assessment") -> str:
        """
        Create risk assessment presentation
        
        Args:
            report_data: Report data
            filename: Output filename
            
        Returns:
            Path to generated presentation
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, report_data, "Risk Assessment")
            
            # Risk overview
            self._create_risk_overview_slide(prs, report_data)
            
            # VaR analysis
            self._create_var_analysis_slide(prs, report_data)
            
            # Stress testing
            self._create_stress_testing_slide(prs, report_data)
            
            # Scenario analysis
            self._create_scenario_analysis_slide(prs, report_data)
            
            # Risk limits monitoring
            self._create_risk_limits_slide(prs, report_data)
            
            # Tail risk analysis
            self._create_tail_risk_slide(prs, report_data)
            
            # Risk attribution
            self._create_risk_attribution_slide(prs, report_data)
            
            # Risk recommendations
            self._create_risk_recommendations_slide(prs, report_data)
            
            # Save presentation
            ppt_path = self.output_dir / f"{filename}.pptx"
            prs.save(str(ppt_path))
            
            logger.info(f"Risk assessment presentation created: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error creating risk assessment presentation: {e}")
            return ""
    
    def create_strategy_comparison_presentation(self,
                                             report_data: Dict[str, Any],
                                             filename: str = "strategy_comparison") -> str:
        """
        Create strategy comparison presentation
        
        Args:
            report_data: Report data
            filename: Output filename
            
        Returns:
            Path to generated presentation
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, report_data, "Strategy Comparison")
            
            # Comparison overview
            self._create_comparison_overview_slide(prs, report_data)
            
            # Performance comparison
            self._create_performance_comparison_slide(prs, report_data)
            
            # Risk comparison
            self._create_risk_comparison_slide(prs, report_data)
            
            # Risk-return analysis
            self._create_risk_return_analysis_slide(prs, report_data)
            
            # Correlation analysis
            self._create_correlation_analysis_slide(prs, report_data)
            
            # Rankings and scores
            self._create_rankings_slide(prs, report_data)
            
            # Diversification benefits
            self._create_diversification_slide(prs, report_data)
            
            # Strategy recommendations
            self._create_strategy_recommendations_slide(prs, report_data)
            
            # Save presentation
            ppt_path = self.output_dir / f"{filename}.pptx"
            prs.save(str(ppt_path))
            
            logger.info(f"Strategy comparison presentation created: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error creating strategy comparison presentation: {e}")
            return ""
    
    def create_regulatory_report_presentation(self,
                                           report_data: Dict[str, Any],
                                           filename: str = "regulatory_report") -> str:
        """
        Create regulatory report presentation
        
        Args:
            report_data: Report data
            filename: Output filename
            
        Returns:
            Path to generated presentation
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._create_title_slide(prs, report_data, "Regulatory Report")
            
            # Compliance summary
            self._create_compliance_summary_slide(prs, report_data)
            
            # Risk metrics compliance
            self._create_risk_metrics_compliance_slide(prs, report_data)
            
            # Position limits monitoring
            self._create_position_limits_slide(prs, report_data)
            
            # Liquidity risk assessment
            self._create_liquidity_risk_slide(prs, report_data)
            
            # Operational risk
            self._create_operational_risk_slide(prs, report_data)
            
            # Model validation
            self._create_model_validation_slide(prs, report_data)
            
            # Audit trail
            self._create_audit_trail_slide(prs, report_data)
            
            # Save presentation
            ppt_path = self.output_dir / f"{filename}.pptx"
            prs.save(str(ppt_path))
            
            logger.info(f"Regulatory report presentation created: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error creating regulatory report presentation: {e}")
            return ""
    
    def _create_title_slide(self, prs: Presentation, report_data: Dict[str, Any], subtitle: str = "Performance Report"):
        """Create title slide"""
        try:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = f"{self.config.company_name}"
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Subtitle
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_shape.text = f"{subtitle}\n{report_data.get('report_metadata', {}).get('strategy_name', 'Strategy Analysis')}"
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[1].font.size = Pt(18)
            subtitle_shape.text_frame.paragraphs[1].font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
            
            # Add date
            date_text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
            text_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.text = date_text
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['dark'])
            
        except Exception as e:
            logger.error(f"Error creating title slide: {e}")
    
    def _create_executive_summary_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create executive summary slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Executive Summary"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Content
            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            # Overall rating
            exec_summary = report_data.get('executive_summary', {})
            if exec_summary:
                p = text_frame.paragraphs[0]
                p.text = f"Overall Rating: {exec_summary.get('overall_rating', 'N/A')}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
                
                # Performance score
                p = text_frame.add_paragraph()
                p.text = f"Performance Score: {exec_summary.get('performance_score', 0):.1f}/100"
                p.font.size = Pt(16)
                p.level = 1
                
                # Risk score
                p = text_frame.add_paragraph()
                p.text = f"Risk Score: {exec_summary.get('risk_score', 0):.1f}/100"
                p.font.size = Pt(16)
                p.level = 1
                
                # Key insights
                key_insights = exec_summary.get('key_insights', [])
                if key_insights:
                    p = text_frame.add_paragraph()
                    p.text = "Key Insights:"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
                    
                    for insight in key_insights[:3]:  # Top 3 insights
                        p = text_frame.add_paragraph()
                        p.text = f"• {insight}"
                        p.font.size = Pt(14)
                        p.level = 1
            
        except Exception as e:
            logger.error(f"Error creating executive summary slide: {e}")
    
    def _create_key_metrics_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create key metrics slide"""
        try:
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "Key Performance Metrics"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Metrics grid
            key_metrics = report_data.get('key_metrics', {})
            if key_metrics:
                metrics_list = list(key_metrics.items())
                
                # Create 2x3 grid of metrics
                for i, (key, value) in enumerate(metrics_list[:6]):
                    row = i // 3
                    col = i % 3
                    
                    # Position calculation
                    x = Inches(0.5 + col * 4.25)
                    y = Inches(2 + row * 2.5)
                    width = Inches(4)
                    height = Inches(2)
                    
                    # Create metric box
                    metric_box = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, x, y, width, height
                    )
                    
                    # Format metric box
                    fill = metric_box.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(self.config.brand_colors['light'])
                    
                    line = metric_box.line
                    line.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
                    line.width = Pt(2)
                    
                    # Add metric text
                    text_frame = metric_box.text_frame
                    text_frame.clear()
                    
                    # Metric value
                    p = text_frame.paragraphs[0]
                    if isinstance(value, (int, float)):
                        if 'return' in key.lower() or 'ratio' in key.lower():
                            p.text = f"{value:.2%}"
                        else:
                            p.text = f"{value:.4f}"
                    else:
                        p.text = str(value)
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
                    p.alignment = PP_ALIGN.CENTER
                    
                    # Metric label
                    p = text_frame.add_paragraph()
                    p.text = key.replace('_', ' ').title()
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['dark'])
                    p.alignment = PP_ALIGN.CENTER
            
        except Exception as e:
            logger.error(f"Error creating key metrics slide: {e}")
    
    def _create_performance_overview_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create performance overview slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Performance Overview"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Create performance chart placeholder
            # In a real implementation, you would generate the chart and embed it
            chart_placeholder = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
            chart_frame = chart_placeholder.text_frame
            chart_frame.text = "Performance Chart\n(Cumulative Returns, Drawdown, Risk Metrics)"
            chart_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            chart_frame.paragraphs[0].font.size = Pt(18)
            chart_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
            
            # Add border to chart placeholder
            shape = chart_placeholder
            line = shape.line
            line.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            line.width = Pt(2)
            
        except Exception as e:
            logger.error(f"Error creating performance overview slide: {e}")
    
    def _create_risk_assessment_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create risk assessment slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Risk Assessment"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Risk metrics content
            risk_assessment = report_data.get('risk_assessment', {})
            if risk_assessment:
                content = slide.shapes.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                # Overall risk score
                p = text_frame.paragraphs[0]
                p.text = f"Overall Risk Score: {risk_assessment.get('risk_score', 0):.1f}/100"
                p.font.size = Pt(20)
                p.font.bold = True
                
                # Risk rating
                p = text_frame.add_paragraph()
                p.text = f"Risk Rating: {risk_assessment.get('risk_rating', 'N/A')}"
                p.font.size = Pt(18)
                p.level = 1
                
                # VaR metrics
                p = text_frame.add_paragraph()
                p.text = f"VaR (95%): {risk_assessment.get('var_95', 0):.2%}"
                p.font.size = Pt(16)
                p.level = 1
                
                # Expected shortfall
                p = text_frame.add_paragraph()
                p.text = f"Expected Shortfall: {risk_assessment.get('expected_shortfall', 0):.2%}"
                p.font.size = Pt(16)
                p.level = 1
                
                # Tail ratio
                p = text_frame.add_paragraph()
                p.text = f"Tail Ratio: {risk_assessment.get('tail_ratio', 0):.2f}"
                p.font.size = Pt(16)
                p.level = 1
            
        except Exception as e:
            logger.error(f"Error creating risk assessment slide: {e}")
    
    def _create_key_insights_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create key insights slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Key Insights"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Insights content
            key_insights = report_data.get('key_insights', [])
            if key_insights:
                content = slide.shapes.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                for i, insight in enumerate(key_insights):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = f"• {insight}"
                    p.font.size = Pt(16)
                    p.level = 0
                    
                    # Color alternate insights
                    if i % 2 == 0:
                        p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['dark'])
                    else:
                        p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
            
        except Exception as e:
            logger.error(f"Error creating key insights slide: {e}")
    
    def _create_recommendations_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create recommendations slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Recommendations"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Recommendations content
            recommendations = report_data.get('recommendations', [])
            if recommendations:
                content = slide.shapes.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()
                
                for i, recommendation in enumerate(recommendations):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = f"{i+1}. {recommendation}"
                    p.font.size = Pt(16)
                    p.level = 0
                    p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['dark'])
            
        except Exception as e:
            logger.error(f"Error creating recommendations slide: {e}")
    
    def _create_appendix_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create appendix slide"""
        try:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = "Appendix"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Appendix content
            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            # Methodology
            p = text_frame.paragraphs[0]
            p.text = "Methodology & Assumptions"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
            
            # Add methodology details
            p = text_frame.add_paragraph()
            p.text = "• Performance calculations based on daily returns"
            p.font.size = Pt(14)
            p.level = 1
            
            p = text_frame.add_paragraph()
            p.text = "• Risk metrics calculated using 252-day trading year"
            p.font.size = Pt(14)
            p.level = 1
            
            p = text_frame.add_paragraph()
            p.text = "• VaR calculated at 95% confidence level"
            p.font.size = Pt(14)
            p.level = 1
            
            # Contact information
            p = text_frame.add_paragraph()
            p.text = "Contact Information"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(self.config.brand_colors['secondary'])
            
            p = text_frame.add_paragraph()
            p.text = f"Generated by {self.config.company_name}"
            p.font.size = Pt(14)
            p.level = 1
            
            p = text_frame.add_paragraph()
            p.text = f"Report Date: {datetime.now().strftime('%Y-%m-%d')}"
            p.font.size = Pt(14)
            p.level = 1
            
        except Exception as e:
            logger.error(f"Error creating appendix slide: {e}")
    
    # Additional slide creation methods would continue here...
    # (Implementation of other slide types)
    
    def _create_performance_summary_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create performance summary slide"""
        try:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "Performance Summary"
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.config.brand_colors['primary'])
            
            # Add performance summary content
            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            
            # This would typically include performance analysis data
            p = text_frame.paragraphs[0]
            p.text = "Performance Summary will be populated with actual data"
            p.font.size = Pt(16)
            
        except Exception as e:
            logger.error(f"Error creating performance summary slide: {e}")
    
    def _create_returns_analysis_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create returns analysis slide"""
        # Similar implementation pattern as above
        pass
    
    def _create_risk_adjusted_metrics_slide(self, prs: Presentation, report_data: Dict[str, Any]):
        """Create risk-adjusted metrics slide"""
        # Similar implementation pattern as above
        pass
    
    # Additional helper methods for other slide types...
    
    def _load_templates(self):
        """Load presentation templates"""
        try:
            # Load HTML email templates
            self.templates['email_summary'] = self._get_email_template('summary')
            self.templates['email_detailed'] = self._get_email_template('detailed')
            self.templates['email_alert'] = self._get_email_template('alert')
            
            logger.info("Presentation templates loaded")
            
        except Exception as e:
            logger.error(f"Error loading templates: {e}")
    
    def _get_email_template(self, template_type: str) -> Template:
        """Get email template"""
        
        if template_type == 'summary':
            template_str = """
            <html>
            <head>
                <style>
                    body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
                    .header { background-color: #2E86AB; color: white; padding: 20px; text-align: center; }
                    .content { padding: 20px; }
                    .metric { display: inline-block; margin: 10px; padding: 15px; background-color: #f8f9fa; border-radius: 5px; }
                    .footer { background-color: #f8f9fa; padding: 20px; text-align: center; margin-top: 30px; }
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>{{ company_name }}</h1>
                    <h2>Performance Summary</h2>
                </div>
                <div class="content">
                    <h3>Key Metrics</h3>
                    {% for key, value in metrics.items() %}
                    <div class="metric">
                        <strong>{{ key }}:</strong> {{ value }}
                    </div>
                    {% endfor %}
                </div>
                <div class="footer">
                    <p>Generated on {{ timestamp }}</p>
                </div>
            </body>
            </html>
            """
        elif template_type == 'detailed':
            template_str = """
            <html>
            <head>
                <style>
                    body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
                    .header { background-color: #2E86AB; color: white; padding: 20px; text-align: center; }
                    .section { margin: 20px 0; padding: 15px; border-left: 4px solid #2E86AB; }
                    .table { width: 100%; border-collapse: collapse; margin: 15px 0; }
                    .table th, .table td { border: 1px solid #ddd; padding: 8px; text-align: left; }
                    .table th { background-color: #f2f2f2; }
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>{{ company_name }}</h1>
                    <h2>Detailed Performance Report</h2>
                </div>
                <div class="section">
                    <h3>Executive Summary</h3>
                    <p>{{ executive_summary }}</p>
                </div>
                <div class="section">
                    <h3>Performance Metrics</h3>
                    <table class="table">
                        <tr><th>Metric</th><th>Value</th></tr>
                        {% for key, value in metrics.items() %}
                        <tr><td>{{ key }}</td><td>{{ value }}</td></tr>
                        {% endfor %}
                    </table>
                </div>
            </body>
            </html>
            """
        elif template_type == 'alert':
            template_str = """
            <html>
            <head>
                <style>
                    body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
                    .header { background-color: #DC3545; color: white; padding: 20px; text-align: center; }
                    .alert { background-color: #f8d7da; border: 1px solid #f5c6cb; padding: 15px; margin: 20px 0; border-radius: 5px; }
                    .warning { background-color: #fff3cd; border: 1px solid #ffeaa7; }
                    .info { background-color: #d1ecf1; border: 1px solid #bee5eb; }
                </style>
            </head>
            <body>
                <div class="header">
                    <h1>{{ company_name }}</h1>
                    <h2>Alert Notification</h2>
                </div>
                <div class="alert {{ alert_type }}">
                    <h3>{{ alert_title }}</h3>
                    <p>{{ alert_message }}</p>
                    <p><strong>Timestamp:</strong> {{ timestamp }}</p>
                </div>
            </body>
            </html>
            """
        else:
            template_str = "<html><body><h1>Template not found</h1></body></html>"
        
        return Template(template_str)
    
    def generate_email_template(self, 
                              template_type: str,
                              data: Dict[str, Any]) -> str:
        """
        Generate email template
        
        Args:
            template_type: Type of email template
            data: Template data
            
        Returns:
            Rendered HTML email
        """
        try:
            template = self.templates.get(f'email_{template_type}')
            if not template:
                return ""
            
            # Add company name and timestamp
            data['company_name'] = self.config.company_name
            data['timestamp'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            
            return template.render(**data)
            
        except Exception as e:
            logger.error(f"Error generating email template: {e}")
            return ""


# Global instance
presentation_templates = PresentationTemplates()